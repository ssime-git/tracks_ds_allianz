from __future__ import annotations

import csv
import random
import shutil
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR


ROOT = Path(__file__).resolve().parents[1]
SLIDES = ROOT / "slides" / "week1_python_fundamentals_databricks.pptx"
DATA = ROOT / "data" / "insurance.csv"
NOTES = ROOT / "slides" / "week1_speaker_notes.md"
LIORA_TEMPLATE = Path(
    "/Users/seb/Downloads/work_cc/slides-mc/liora_claude_design_template/"
    "Jour 1 - Fondations LLM (fixed).pptx"
)

NAVY = RGBColor(15, 32, 56)
INK = RGBColor(31, 41, 55)
MUTED = RGBColor(96, 108, 128)
ORANGE = RGBColor(244, 119, 56)
LIME = RGBColor(183, 222, 69)
BLUE = RGBColor(56, 116, 203)
PALE = RGBColor(246, 248, 252)
WHITE = RGBColor(255, 255, 255)


def add_textbox(slide, x, y, w, h, text, size=20, color=INK, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Inter"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_title(slide, section, title, subtitle=None):
    add_textbox(slide, 0.55, 0.28, 2.4, 0.28, section.upper(), 9, ORANGE, True)
    add_textbox(slide, 0.55, 0.58, 8.2, 0.72, title, 30, NAVY, True)
    if subtitle:
        add_textbox(slide, 0.58, 1.22, 8.3, 0.45, subtitle, 13, MUTED)


def footer(slide, n):
    add_textbox(slide, 11.7, 6.9, 1.0, 0.2, f"{n:02d}", 8, MUTED, True, PP_ALIGN.RIGHT)


def add_band(slide, color=NAVY):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_card(slide, x, y, w, h, title, body, accent=ORANGE):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = RGBColor(225, 231, 240)
    card.adjustments[0] = 0.08
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    add_textbox(slide, x + 0.22, y + 0.18, w - 0.35, 0.35, title, 15, NAVY, True)
    add_textbox(slide, x + 0.22, y + 0.62, w - 0.35, h - 0.75, body, 11, MUTED)


def add_flow(slide, labels, x=1.2, y=2.25, w=2.1, h=0.72, gap=0.65):
    for i, label in enumerate(labels):
        sx = x + i * (w + gap)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(sx), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = WHITE
        shape.line.color.rgb = ORANGE if i == 1 else RGBColor(219, 226, 236)
        add_textbox(slide, sx + 0.1, y + 0.2, w - 0.2, 0.25, label, 13, NAVY, True, PP_ALIGN.CENTER)
        if i < len(labels) - 1:
            line = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                Inches(sx + w + 0.05),
                Inches(y + h / 2),
                Inches(sx + w + gap - 0.1),
                Inches(y + h / 2),
            )
            line.line.color.rgb = MUTED
            line.line.width = Pt(1.4)


def add_table_like(slide, x, y, headers, rows, col_w, row_h=0.43):
    for c, head in enumerate(headers):
        rx = x + sum(col_w[:c])
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(rx), Inches(y), Inches(col_w[c]), Inches(row_h))
        rect.fill.solid()
        rect.fill.fore_color.rgb = NAVY
        rect.line.color.rgb = WHITE
        add_textbox(slide, rx + 0.06, y + 0.11, col_w[c] - 0.12, 0.16, head, 9, WHITE, True)
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            rx = x + sum(col_w[:c])
            ry = y + row_h * (r + 1)
            rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(rx), Inches(ry), Inches(col_w[c]), Inches(row_h))
            rect.fill.solid()
            rect.fill.fore_color.rgb = PALE if r % 2 == 0 else WHITE
            rect.line.color.rgb = RGBColor(230, 234, 241)
            add_textbox(slide, rx + 0.06, ry + 0.11, col_w[c] - 0.12, 0.16, str(val), 9, INK)


def code_block(slide, x, y, w, h, code):
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(26, 33, 46)
    rect.line.color.rgb = RGBColor(26, 33, 46)
    add_textbox(slide, x + 0.2, y + 0.18, w - 0.4, h - 0.25, code, 15, WHITE)


def make_dataset():
    random.seed(42)
    regions = ["northeast", "northwest", "southeast", "southwest"]
    sexes = ["female", "male"]
    DATA.parent.mkdir(parents=True, exist_ok=True)
    with DATA.open("w", newline="") as f:
        writer = csv.writer(f)
        writer.writerow(["age", "sex", "bmi", "children", "smoker", "region", "charges"])
        for _ in range(420):
            age = random.randint(18, 64)
            sex = random.choice(sexes)
            bmi = round(max(16, min(45, random.gauss(30, 6))), 1)
            children = random.choices([0, 1, 2, 3, 4, 5], [0.42, 0.2, 0.18, 0.12, 0.05, 0.03])[0]
            smoker = "yes" if random.random() < 0.2 else "no"
            region = random.choice(regions)
            base = 1800 + age * 190 + bmi * 155 + children * 520
            if smoker == "yes":
                base += 18500 + bmi * 210
            if region == "southeast":
                base += 900
            noise = random.gauss(0, 2200)
            charges = round(max(1200, base + noise), 2)
            writer.writerow([age, sex, bmi, children, smoker, region, charges])


def slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def make_deck_custom():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    notes = []

    specs = [
        ("cover", "FORMATION · DATA SCIENTIST TRACK", "Python Fundamentals\nin Azure Databricks", "Semaine 1 · SAS vers Python · pandas · Spark · SQL"),
        ("cards", "Objectifs", "À la fin de la semaine 1, vous savez...", [
            ("Naviguer", "Ouvrir un notebook, attacher un cluster, exécuter une cellule."),
            ("Coder simple", "Lire variables, types, listes, fonctions et conditions."),
            ("Analyser", "Charger, filtrer et agréger une table avec pandas."),
            ("Traduire", "Relier DATA step, PROC SQL et DataFrame Databricks."),
        ]),
        ("flow", "Pourquoi Databricks?", "Un espace commun pour exécuter, documenter et partager l'analyse.", ["Data sources", "Notebook", "Python / SQL / Spark", "Analytics / ML"]),
        ("workspace", "Workspace Azure Databricks", "Les 5 zones à reconnaître avant de coder."),
        ("notebook", "Le notebook", "Code, résultat et explication dans le même support."),
        ("table", "SAS vs Databricks", "Même logique analytique, nouvelle interface.", ["SAS", "Databricks"], [
            ["SAS Program", "Notebook"],
            ["DATA Step", "pandas / Spark"],
            ["PROC SQL", "SQL / Spark SQL"],
            ["SAS Dataset", "DataFrame"],
            ["SAS Libraries", "Catalog / Storage"],
        ]),
        ("code", "Variables", "Une variable donne un nom à une valeur.", 'age = 45\nregion = "North"\npremium = 1200.50'),
        ("table", "Types de données", "Les 4 types à reconnaître le premier jour.", ["Type", "Exemple"], [["int", "45"], ["float", "12.5"], ["str", '"Paris"'], ["bool", "True"]]),
        ("code", "Listes", "Une liste regroupe plusieurs valeurs.", 'regions = ["North", "South", "East"]\nregions[0]'),
        ("code", "Fonctions", "Une fonction rend une règle de calcul réutilisable.", "def add_tax(x):\n    return x * 1.2\n\nadd_tax(1000)"),
        ("indent", "Lire du Python", "L'indentation montre ce qui appartient au bloc."),
        ("flow", "Workflow analytics Python", "Un enchaînement simple, proche des habitudes SAS.", ["Read Data", "Transform", "Analyze", "Visualize"]),
        ("dataframe", "Qu'est-ce qu'un DataFrame?", "Une table: lignes, colonnes, types et opérations."),
        ("code", "Charger un CSV", "pandas lit le fichier et crée un DataFrame.", 'import pandas as pd\n\ndf = pd.read_csv("/dbfs/FileStore/tables/insurance.csv")'),
        ("code", "Première exploration", "Trois réflexes avant toute analyse.", "df.head()\ndf.info()\ndf.describe()"),
        ("code", "Filtrer des lignes", "On garde les lignes qui respectent une condition.", 'df[df["age"] > 50]'),
        ("code", "Sélectionner des colonnes", "On réduit la table aux variables utiles.", 'df[["age", "charges"]]'),
        ("groupby", "GroupBy", "L'équivalent mental d'une agrégation par classe.", 'df.groupby("region")["charges"].mean()'),
        ("table", "SAS vs pandas", "Repères pour traduire les gestes connus.", ["SAS", "pandas"], [["PROC FREQ", "value_counts()"], ["PROC MEANS", "groupby() + agg()"], ["DATA Step", "assign() / filters"], ["WHERE", "boolean filter"]]),
        ("flow", "Mini workflow", "Question métier: charges moyennes par statut fumeur.", ["Load", "Filter", "Aggregate", "Interpret"]),
        ("sparkwhy", "Pourquoi Spark?", "Spark garde le modèle table, mais vise des volumes plus grands."),
        ("code", "Lire avec Spark", "Le résultat est un Spark DataFrame.", 'spark_df = spark.read.csv(\n    "/FileStore/tables/insurance.csv",\n    header=True,\n    inferSchema=True\n)'),
        ("code", "Transformations Spark", "Même intention: choisir et filtrer.", 'spark_df.select("age", "charges")\nspark_df.filter(spark_df.age > 50)'),
        ("flow", "Architecture simplifiée", "Pour cette semaine, retenir seulement ces trois briques.", ["Notebook", "Cluster", "Data"]),
        ("code", "Cellules SQL", "SQL reste disponible dans le notebook.", "SELECT region,\n       AVG(charges) AS avg_charges\nFROM insurance\nGROUP BY region"),
        ("cards", "Collaboration", "Databricks sert aussi à travailler ensemble.", [
            ("Partage", "Notebooks dans un dossier commun."),
            ("Exécution", "Cluster partagé pour la session."),
            ("Traçabilité", "Code, commentaires et résultats au même endroit."),
            ("Réutilisation", "Même notebook pour démo et exercices."),
        ]),
        ("cards", "Exercices guidés", "Chaque exercice part d'un code existant.", [
            ("Explorer", "head, columns, describe."),
            ("Filtrer", "age > 50, smoker, BMI."),
            ("Agréger", "charges moyennes par segment."),
            ("Traduire", "pandas vers Spark et SQL."),
        ]),
        ("challenge", "Mini challenge final", "Quel segment génère les coûts médicaux les plus élevés?"),
        ("summary", "Synthèse", "Les concepts à retenir avant la semaine 2."),
    ]

    for idx, spec in enumerate(specs, 1):
        s = slide(prs)
        kind = spec[0]
        if kind == "cover":
            add_band(s, NAVY)
            add_textbox(s, 0.75, 0.65, 5.3, 0.3, spec[1], 10, LIME, True)
            add_textbox(s, 0.75, 1.55, 8.4, 1.65, spec[2], 38, WHITE, True)
            add_textbox(s, 0.8, 3.55, 7.4, 0.55, spec[3], 16, RGBColor(221, 231, 244))
            add_flow(s, ["SAS", "Python", "pandas", "Spark", "SQL"], x=0.85, y=5.35, w=1.65, h=0.62, gap=0.34)
        elif kind == "cards":
            add_title(s, "Week 1", spec[1], spec[2])
            for i, (t, b) in enumerate(spec[3]):
                add_card(s, 0.75 + (i % 2) * 6.0, 2.0 + (i // 2) * 1.75, 5.25, 1.25, t, b, [ORANGE, LIME, BLUE, ORANGE][i])
        elif kind == "flow":
            add_title(s, "Concept", spec[1], spec[2])
            add_flow(s, spec[3], x=0.95, y=3.1, w=2.35, h=0.82, gap=0.45)
        elif kind == "workspace":
            add_title(s, "Databricks", spec[1], spec[2])
            add_card(s, 0.8, 2.0, 2.1, 3.7, "Workspace", "Dossiers partagés\nNotebooks\nPermissions", ORANGE)
            add_card(s, 3.25, 2.0, 2.1, 3.7, "Compute", "Cluster\nRun\nAuto-terminate", BLUE)
            add_card(s, 5.7, 2.0, 2.1, 3.7, "Catalog", "Tables\nSchémas\nVues", LIME)
            add_card(s, 8.15, 2.0, 2.1, 3.7, "Data", "CSV\nStorage\nFileStore", ORANGE)
            add_card(s, 10.6, 2.0, 2.1, 3.7, "Notebook", "Markdown\nPython\nSQL", BLUE)
        elif kind == "notebook":
            add_title(s, "Notebook", spec[1], spec[2])
            add_card(s, 1.0, 2.0, 3.2, 2.6, "Markdown", "Texte explicatif\nConsignes\nInterprétation", ORANGE)
            code_block(s, 4.65, 2.0, 3.4, 2.6, 'age = 45\npremium = 1000\npremium * 1.2')
            add_card(s, 8.5, 2.0, 3.2, 2.6, "SQL", "SELECT region,\n       AVG(charges)\nFROM insurance", LIME)
        elif kind == "table":
            add_title(s, "Mapping", spec[1], spec[2])
            add_table_like(s, 2.0, 2.15, spec[3], spec[4], [3.2, 4.1])
        elif kind == "code":
            add_title(s, "Python", spec[1], spec[2])
            code_block(s, 1.35, 2.25, 6.2, 2.75, spec[3])
            add_card(s, 8.2, 2.25, 3.8, 2.75, "Lecture guidée", "1. Lire de gauche à droite.\n2. Identifier les noms.\n3. Repérer les symboles.\n4. Exécuter puis observer.", LIME)
        elif kind == "indent":
            add_title(s, "Python", spec[1], spec[2])
            code_block(s, 1.3, 2.0, 5.3, 2.9, 'if smoker == "yes":\n    print("Smoker")\n    print("Higher risk")\n\nprint("Always executed")')
            add_card(s, 7.15, 2.0, 4.35, 2.9, "Point clé", "Les lignes décalées appartiennent au bloc.\n\nOn lit lentement: condition, deux actions, puis retour au flux normal.", ORANGE)
        elif kind == "dataframe":
            add_title(s, "pandas", spec[1], spec[2])
            add_table_like(s, 1.2, 2.0, ["age", "sex", "bmi", "smoker", "charges"], [[45, "male", 28.4, "no", 8231], [52, "female", 31.2, "yes", 28910], [29, "female", 24.1, "no", 4120]], [1.2, 1.3, 1.2, 1.5, 1.6])
            add_card(s, 8.5, 2.0, 3.2, 2.15, "Mental model", "Excel-like table\nSAS dataset\nDataFrame pandas", LIME)
        elif kind == "groupby":
            add_title(s, "pandas", spec[1], spec[2])
            code_block(s, 0.95, 2.0, 5.3, 1.5, spec[3])
            add_flow(s, ["Rows", "Group", "Mean", "Result"], x=1.0, y=4.35, w=1.8, h=0.62, gap=0.5)
            add_card(s, 7.35, 2.0, 4.0, 2.7, "À dire", "On ne change pas la question métier.\nOn change seulement l'outil utilisé pour la poser.", ORANGE)
        elif kind == "sparkwhy":
            add_title(s, "Spark", spec[1], spec[2])
            add_card(s, 1.2, 2.25, 4.4, 2.2, "pandas", "Très pratique pour apprendre et explorer.\nTravaille dans la mémoire de la machine.", ORANGE)
            add_card(s, 7.1, 2.25, 4.4, 2.2, "Spark", "Même logique de table.\nPrévu pour traiter plus gros dans Databricks.", BLUE)
        elif kind == "challenge":
            add_title(s, "Exercice", spec[1], spec[2])
            add_flow(s, ["Choose segment", "Group", "Average charges", "Interpret"], x=1.0, y=2.4, w=2.35, h=0.78, gap=0.42)
            add_card(s, 2.4, 4.25, 8.2, 1.2, "Livrable attendu", "Une table triée et une phrase métier: le segment le plus coûteux est ... parce que ...", LIME)
        elif kind == "summary":
            add_title(s, "Wrap-up", spec[1], spec[2])
            items = [("Notebook", "Exécuter et documenter"), ("Python", "Lire les bases"), ("pandas", "Manipuler une table"), ("Spark", "Découvrir les gros volumes"), ("SQL", "Réutiliser les acquis")]
            for i, (t, b) in enumerate(items):
                add_card(s, 0.75 + i * 2.45, 2.3, 2.05, 2.35, t, b, [ORANGE, BLUE, LIME, ORANGE, BLUE][i])
        footer(s, idx)
        notes.append(f"## Slide {idx:02d} - {spec[1] if kind != 'cover' else 'Cover'}\n\n"
                     "- Parler lentement et relier au vocabulaire SAS.\n"
                     "- Faire exécuter ou lire un exemple concret avant de passer à la suite.\n")

    SLIDES.parent.mkdir(parents=True, exist_ok=True)
    prs.save(SLIDES)
    NOTES.write_text("\n".join(notes), encoding="utf-8")


def _delete_slide(prs, index):
    slide_id_list = prs.slides._sldIdLst
    slides = list(slide_id_list)
    slide_id_list.remove(slides[index])


def _text_shapes(slide):
    return [
        shape
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False) and shape.text_frame is not None
    ]


def _set_shape_text(shape, text):
    tf = shape.text_frame
    first_p = tf.paragraphs[0]
    first_run = first_p.runs[0] if first_p.runs else first_p.add_run()
    first_run.text = text
    for run in first_p.runs[1:]:
        run.text = ""
    for p in tf.paragraphs[1:]:
        for run in p.runs:
            run.text = ""


def _replace_texts(slide, replacements):
    shapes = _text_shapes(slide)
    for shape, text in zip(shapes, replacements):
        _set_shape_text(shape, text)
    for shape in shapes[len(replacements):]:
        _set_shape_text(shape, "")


def make_deck():
    """Generate the Week 1 deck by reusing the actual Liora PowerPoint template."""
    SLIDES.parent.mkdir(parents=True, exist_ok=True)
    shutil.copyfile(LIORA_TEMPLATE, SLIDES)
    prs = Presentation(SLIDES)

    while len(prs.slides) > 29:
        _delete_slide(prs, len(prs.slides) - 1)

    slide_texts = [
        [
            "Formation · Data Scientist Track",
            "Python",
            "Fundamentals",
            "in Azure Databricks.",
            "SAS vers Python, pandas, Spark et SQL",
            "Week 01 / 06",
            "Week 1 · Databricks",
            "01",
        ],
        [
            "Avant de commencer",
            "Cadre de la semaine 1",
            "01 · Posture",
            "Analyste, pas développeur",
            "On apprend à lire, exécuter et adapter du code analytique simple.",
            "02 · Progression",
            "Très guidée",
            "Aucun exercice page blanche. On modifie du code existant.",
            "03 · Échanges",
            "Questions à tout moment",
            "On relie chaque geste au vocabulaire SAS connu.",
            "04 · Données",
            "Dataset pédagogique",
            "Fichier insurance.csv: âge, BMI, fumeur, région, charges.",
        ],
        [
            "Cas fil rouge",
            "Analyser les coûts médicaux avec Databricks.",
            "À la fin de la semaine 1, vous aurez exécuté un workflow complet :",
            "Charger → un CSV insurance dans Databricks.",
            "Explorer → lignes, colonnes, types et statistiques.",
            "Transformer → filtrer, sélectionner, agréger.",
            "Comparer → pandas, Spark et SQL sur la même question.",
            "[Input] insurance.csv",
            "[Notebook] Python + Markdown",
            "[DataFrame] pandas puis Spark",
            "[SQL] Requêtes analytiques",
            "[Output] Segment de coût élevé",
        ],
        [
            "Programme",
            "Une brique par séquence, un notebook par geste clé.",
            "Chaque séquence démarre par une démonstration lente, puis un exercice guidé transposable.",
            "Chaque séquence démarre par une démonstration lente, puis un exercice guidé transposable.",
            "Programme",
            "04",
        ],
        [
            "Au programme",
            "Sommaire de la session live",
            "01 · Databricks → workspace, cluster, notebook — 35 min",
            "02 · Python basics → variables, types, listes, fonctions — 45 min",
            "03 · Lire du code → indentation, conditions, erreurs fréquentes — 25 min",
            "04 · pandas → charger, explorer, filtrer, grouper — 55 min",
            "05 · Spark → DataFrame, select, filter, groupBy — 30 min",
            "06 · SQL → cellules SQL, COUNT, AVG, GROUP BY — 35 min",
            "07 · Mini challenge → segment le plus coûteux — 25 min",
            "08 · Synthèse → checklist Databricks et SAS mapping — 15 min",
            "Sommaire",
            "05",
        ],
        ["01", "Section 01 · 35 min", "Découvrir Databricks.", "On apprend à se repérer avant d'écrire du code.", "06 Section 01 — Databricks", "06"],
        [
            "01 · Databricks",
            "À la fin de la section, vous savez…",
            "01",
            "Ouvrir un notebook",
            "Retrouver un support dans /Shared/training et lire les cellules.",
            "02",
            "Attacher un cluster",
            "Comprendre pourquoi le notebook a besoin d'un compute.",
            "03",
            "Exécuter une cellule",
            "Lancer du Markdown, du Python ou du SQL et lire le résultat.",
            "04",
            "Repérer la donnée",
            "Comprendre le chemin FileStore utilisé par pandas et Spark.",
            "05",
            "Faire le lien SAS",
            "Programme SAS, DATA step, PROC SQL deviennent notebook, DataFrame, SQL.",
        ],
        [
            "01 · Databricks",
            "Les 4 règles de confort",
            "TABLEAU :",
            "Règle  Pourquoi  Équivalent SAS",
            "Toujours vérifier le cluster  Sans cluster attaché, la cellule ne s'exécute pas.  Session SAS disponible",
            "Exécuter dans l'ordre  Une variable créée plus bas n'existe pas encore.  Ordre des étapes DATA/PROC",
            "Lire le résultat de cellule  Databricks affiche la sortie juste sous le code.  Log + output SAS",
            "Sauver dans le dossier partagé  Les notebooks doivent être retrouvables par le groupe.  Bibliothèque partagée",
            "01 · Databricks",
            "08",
        ],
        ["02", "Section 02 · 45 min", "Python,fondamentaux.", "Des bases suffisantes pour lire et modifier des notebooks analytiques.", "09 Section 02 — Python", "09"],
        [
            "02 · Python",
            "Une variable, en une phrase.",
            "Un nom qui pointe vers une valeur réutilisable.",
            "Pas une colonne de table. Pas une macro SAS. Un objet simple en mémoire.",
            "→ Excellent pour stocker un seuil, une modalité, un chemin de fichier.",
            "→ À éviter pour représenter une table complète.",
            "On lit d'abord le nom, puis la valeur affectée après le signe égal.",
            'age = 45\nregion = "southwest"\npremium = 1200.50',
            "On lit d'abord le nom, puis la valeur affectée après le signe égal.",
            "02 · Python",
            "10",
        ],
        [
            "02 · Python",
            "4 types à reconnaître",
            "Type 1 · int",
            "Entier",
            "45, 64, 0. Sert aux âges, compteurs, nombres d'enfants.",
            "Calculs simples",
            "Type 2 · float",
            "Décimal",
            "27.5, 1200.50. Sert au BMI, aux montants, aux moyennes.",
            "Analyse",
            "Type 3 · str",
            "Texte",
            "southwest, yes, female. Les guillemets indiquent une chaîne.",
            "Modalités",
            "Type 4 · bool",
            "Vrai / faux",
            "True ou False. Résultat naturel d'une condition.",
            "Filtrage",
        ],
        [
            "02 · Python",
            "L'indentation — le piège utile",
            "Pratique analytique",
            "On lit la structure visuellement :",
            "❌ Mauvaise lecture — Les espaces sont décoratifs — Toutes les lignes sont au même niveau.",
            "✓ Bonne lecture — Les lignes décalées appartiennent au bloc — La condition contrôle seulement ce bloc.",
            'if smoker == "yes":\n    print("Smoker segment")\n    print("Higher risk")',
            "• Lire la condition",
            "• Repérer les lignes indentées",
            "• Vérifier ce qui s'exécute toujours",
            "• Exécuter puis observer",
            "02 · Python",
        ],
        [
            "02 · Python",
            "La fonction évite de recopier une règle.",
            "Une fonction prend une entrée et renvoie une sortie.",
            "Pour un analyste, c'est une petite règle métier nommée. Exemple: appliquer une taxe, calculer un coût annuel.",
            "Réutilisable, testable, lisible — et beaucoup moins fragile qu'un copier-coller.",
            "[Entrée] monthly_cost = 120",
            "[Fonction] annual_cost(x)",
            "[Sortie] 1440",
            "EXEMPLE — def annual_cost(x): return x * 12",
            "Réutilisable, testable, lisible — et beaucoup moins fragile qu'un copier-coller.",
            "02 · Python",
            "13",
        ],
        [
            "02 · Python",
            "Les deux exigences non négociables",
            "On ne cherche pas à devenir développeur. On cherche à produire une analyse fiable.",
            "Exigence 01",
            "Lisibilité",
            "Le code doit pouvoir être relu par un autre analyste: noms clairs, étapes courtes, commentaires utiles.",
            "Exigence 02",
            "Reproductibilité",
            "La même donnée et le même notebook doivent donner le même résultat, sans manipulation cachée.",
            "On ne cherche pas à devenir développeur. On cherche à produire une analyse fiable.",
            "02 · Python",
            "14",
        ],
        ["03", "Section 03 · 55 min", "pandas,pour les tables.", "On retrouve le mental model SAS dataset: lignes, colonnes, transformations.", "15 Section 03 — pandas", "15"],
        [
            "03 · pandas 1",
            "DataFrame · la table centrale",
            "Un DataFrame pandas ressemble à une table Excel ou un dataset SAS.",
            "La question métier ne change pas. On change seulement la syntaxe.",
            "TABLEAU :",
            "Objet  Rôle  Réflexe SAS",
            "df  Table en mémoire  Dataset SAS",
            "df.head()  Voir les premières lignes  PROC PRINT(obs=5)",
            "df.info()  Voir colonnes et types  PROC CONTENTS",
            "df.describe()  Statistiques rapides  PROC MEANS",
            "Dataset SAS → DataFrame pandas pour manipuler localement.",
            "La question métier ne change pas. On change seulement la syntaxe.",
        ],
        [
            "03 · pandas 2",
            "Charger, contexte, chemin",
            "Le chemin pandas commence par /dbfs parce que pandas lit via le système de fichiers local Databricks.",
            "Piège",
            "Deux chemins proches",
            "pandas: /dbfs/FileStore/tables/insurance.csv",
            "Spark",
            "Sans /dbfs",
            "Spark: /FileStore/tables/insurance.csv",
            "TABLEAU :",
            "Librairie  Chemin",
            "pandas  /dbfs/FileStore/tables/insurance.csv",
        ],
        [
            "03 · pandas 3, 4, 5",
            "Trois gestes complémentaires",
            "03 · select",
            "Choisir les colonnes",
            'df[["age", "charges"]] réduit la table aux variables utiles pour la question.',
            "04 · filter",
            "Choisir les lignes",
            'df[df["age"] > 50] garde uniquement les lignes qui respectent la condition.',
            "05 · sort",
            "Classer",
            'df.sort_values("charges", ascending=False) montre les coûts les plus élevés.',
            "03 · pandas",
        ],
        [
            "03 · pandas",
            "GroupBy · l'équivalent mental de PROC MEANS",
            "❌ Lecture compliquée — groupby est une magie Python.",
            "✓ Bonne lecture — on découpe la table par région, puis on calcule la moyenne des charges.",
            'df.groupby("region")["charges"].mean()',
            "03 · pandas",
            "19",
        ],
        [
            "03 · pandas",
            "Les 4 dimensions de l'exploration",
            "Une analyse fiable commence par regarder la donnée avant de conclure.",
            "01 · Forme",
            "Lignes, colonnes",
            "Combien d'observations? Quelles variables? Le fichier est-il celui attendu?",
            "02 · Types",
            "Numérique ou texte",
            "age doit être numérique, smoker doit être une modalité.",
            "03 · Valeurs",
            "Min, max, moyenne",
            "Repérer les ordres de grandeur aberrants.",
        ],
        ["04", "Section 04 · 30 min", "Spark,sans théorie.", "Même logique de table, prévu pour des volumes plus grands dans Databricks.", "21 Section 04 — Spark", "21"],
        [
            "04 · Spark",
            "Spark n'est pas une donnée comme une autre.",
            "Un Spark DataFrame reste une table, mais l'exécution est gérée par le cluster.",
            "→ On garde select, filter, groupBy.",
            "→ On évite les détails de distribution cette semaine.",
            "RÉFÉRENCE — Databricks notebook — Spark est disponible directement via la variable spark.",
            "PRATIQUE — display(spark_df) affiche une table interactive dans le notebook.",
            "Notebook Cluster DataFrame select filter groupBy display",
            "04 · Spark",
            "22",
        ],
        [
            "04 · Spark",
            "Pendant la semaine 1 : simplicité.",
            "Le dataset de formation reste petit :",
            "Colonnes → age, sex, bmi, children, smoker, region, charges.",
            "Objectif → comparer pandas, Spark et SQL sur les mêmes questions.",
            "Spark read → spark.read.csv(..., header=True, inferSchema=True).",
            "display → visualisation Databricks.",
            "groupBy → moyenne de charges par segment.",
            '# Extrait\nspark_df = spark.read.csv("/FileStore/tables/insurance.csv", header=True, inferSchema=True)',
            "04 · Spark",
            "23",
        ],
        [
            "04 · Spark",
            "Traduire un geste pandas en Spark",
            "Avant de passer à SQL, on répète les mêmes intentions :",
            "Sources : notebook 04_spark_basics.py, exercices guidés.",
            "01 · Sélection",
            "Choisir des colonnes",
            'spark_df.select("age", "charges")',
            "02 · Filtre",
            "Garder des lignes",
            "spark_df.filter(spark_df.age > 50)",
            "03 · Agrégation",
            "Calculer une moyenne par groupe",
        ],
        [
            "04 · Spark",
            "Ce qui change vraiment",
            "01 · pandas",
            "Exploration locale",
            "Très lisible pour apprendre, manipuler un CSV et raisonner vite.",
            "02 · Spark",
            "Exécution cluster",
            "Même logique table, mais pensée pour Databricks et des volumes plus grands.",
            "03 · SQL",
            "Langage familier",
            "Requêtes déclaratives, souvent proches de PROC SQL.",
            "04 · Spark",
        ],
        ["05", "Section 05 · 35 min", "SQL,dans Databricks.", "On réutilise un langage familier pour poser des questions analytiques.", "27 Section 05 — SQL", "27"],
        [
            "05 · SQL",
            "5 blocs utiles pour une requête analytique",
            "Une requête claire permet de vérifier vite une intuition métier.",
            "Bloc 01 · SELECT",
            "Quelles colonnes ?",
            "region, smoker, AVG(charges)",
            "Bloc 02 · FROM",
            "Quelle table ?",
            "insurance, créée comme vue temporaire Spark.",
            "Bloc 03 · WHERE",
            "Quelles lignes ?",
            "smoker = 'yes' ou age > 50.",
        ],
        [
            "05 · SQL",
            "COUNT · AVG · GROUP BY",
            "01 · COUNT",
            "Combien de lignes ?",
            "SELECT COUNT(*) FROM insurance",
            "→ Vérifie le volume chargé.",
            "02 · AVG",
            "Quelle moyenne ?",
            "SELECT AVG(charges) FROM insurance",
            "→ Donne un ordre de grandeur.",
            "03 · GROUP BY",
            "Par segment",
        ],
        [
            "05 · SQL",
            "Mini challenge final",
            "Question métier : quel segment génère les coûts médicaux les plus élevés ?",
            "01 · Segmenter",
            "Choisir les axes",
            "smoker, age_segment, sex ou region.",
            "02 · Agréger",
            "Calculer AVG(charges)",
            "Comparer les moyennes par segment, pas seulement les lignes extrêmes.",
            "03 · Interpréter",
            "Phrase métier",
            "Le segment le plus coûteux est ... parce que ...",
            "05 · SQL",
            "29",
        ],
    ]

    for idx, texts in enumerate(slide_texts):
        _replace_texts(prs.slides[idx], texts)

    prs.save(SLIDES)
    NOTES.write_text(
        "# Notes formateur - Week 1\n\n"
        "Deck généré depuis le template Liora réel: "
        f"{LIORA_TEMPLATE}\n\n"
        "Fil conducteur: lire lentement le code, exécuter cellule par cellule, "
        "et relier chaque concept à l'équivalent SAS.\n",
        encoding="utf-8",
    )


if __name__ == "__main__":
    make_dataset()
    make_deck()
    print(f"Created {DATA}")
    print(f"Created {SLIDES}")
    print(f"Created {NOTES}")
