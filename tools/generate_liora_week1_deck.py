from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "slides" / "week1_python_fundamentals_databricks.pptx"
NOTES = ROOT / "slides" / "week1_speaker_notes.md"
MAZE = ROOT / ".tmp_liora_zip" / "assets" / "maze.png"

NAVY = RGBColor(26, 26, 51)
CORAL = RGBColor(255, 103, 69)
CREAM = RGBColor(232, 231, 225)
CREAM_SOFT = RGBColor(242, 241, 237)
WHITE = RGBColor(255, 255, 255)
INK_SOFT = RGBColor(96, 96, 116)
PINK = RGBColor(255, 225, 220)
GREEN = RGBColor(197, 239, 206)
BLACK = RGBColor(18, 18, 24)
PURPLE = RGBColor(118, 87, 255)
YELLOW = RGBColor(221, 255, 69)


def rgb(hex_color: str) -> RGBColor:
    hex_color = hex_color.strip("#")
    return RGBColor(int(hex_color[:2], 16), int(hex_color[2:4], 16), int(hex_color[4:], 16))


class Deck:
    def __init__(self) -> None:
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.n = 0
        self.notes: list[str] = []

    def slide(self, bg=CREAM):
        self.n += 1
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, self.prs.slide_height)
        rect.fill.solid()
        rect.fill.fore_color.rgb = bg
        rect.line.fill.background()
        return s

    def text(self, s, x, y, w, h, text, size=24, color=NAVY, bold=False, align=PP_ALIGN.LEFT, font="Arial"):
        box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame
        tf.clear()
        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0
        tf.margin_bottom = 0
        tf.vertical_anchor = MSO_ANCHOR.TOP
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        return box

    def foot(self, s, label):
        self.text(s, 0.55, 7.05, 3.0, 0.16, label.upper(), 6, INK_SOFT, True)
        self.text(s, 12.25, 7.05, 0.5, 0.16, f"{self.n:02d}", 6, INK_SOFT, True, PP_ALIGN.RIGHT)
        self.notes.append(f"{label}: prévoir 2-4 min, ou 8-12 min si la slide lance une démo/exercice.")

    def maze(self, s, where="tr", opacity=False):
        if not MAZE.exists():
            return
        if where == "tr":
            pic = s.shapes.add_picture(str(MAZE), Inches(9.35), Inches(-0.65), width=Inches(4.8))
        elif where == "cover":
            pic = s.shapes.add_picture(str(MAZE), Inches(8.3), Inches(-0.9), width=Inches(6.1))
        else:
            pic = s.shapes.add_picture(str(MAZE), Inches(-1.4), Inches(5.35), width=Inches(4.2))

    def title(self, s, crumb, title, subtitle=None):
        self.text(s, 0.55, 0.45, 2.7, 0.16, crumb.upper(), 6, INK_SOFT, True)
        self.text(s, 0.55, 0.68, 9.4, 0.72, title, 24, NAVY, True)
        if subtitle:
            self.text(s, 0.55, 1.24, 9.7, 0.35, subtitle, 11, INK_SOFT)
        line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(1.72), Inches(11.9), Inches(0.012))
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(205, 202, 193)
        line.line.fill.background()

    def cover(self):
        s = self.slide(NAVY)
        self.maze(s, "cover")
        self.text(s, 0.62, 0.65, 4.2, 0.25, "FORMATION · DATA SCIENTIST TRACK", 9, CORAL, True)
        self.text(s, 0.62, 1.55, 9.4, 2.55, "Python\nFundamentals\nin Databricks", 52, CREAM, True)
        self.text(s, 0.72, 5.15, 4.4, 0.28, "SAS → Python · pandas · Spark · SQL", 13, CREAM)
        pill = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.72), Inches(5.78), Inches(1.7), Inches(0.36))
        pill.fill.solid(); pill.fill.fore_color.rgb = CORAL; pill.line.fill.background()
        self.text(s, 0.92, 5.88, 1.25, 0.12, "WEEK 01 / 06", 7, WHITE, True, PP_ALIGN.CENTER)
        self.text(s, 0.62, 7.04, 2.1, 0.16, "LIORA · ALLIANZ", 6, CREAM, True)
        self.text(s, 12.25, 7.04, 0.5, 0.16, "01", 6, CREAM, True, PP_ALIGN.RIGHT)
        self.notes.append("Cover. Installer le cadre: analystes SAS, pas formation développeur.")

    def divider(self, num, title, lede, label):
        s = self.slide(CORAL)
        self.maze(s, "tr")
        self.text(s, 0.7, 1.15, 2.0, 2.55, f"{num:02d}", 136, NAVY, True)
        self.text(s, 3.45, 1.7, 4.5, 0.2, label.upper(), 8, NAVY, True)
        self.text(s, 3.45, 2.05, 8.3, 1.3, title, 42, NAVY, True)
        self.text(s, 3.48, 3.65, 7.2, 0.7, lede, 15, NAVY)
        self.foot(s, f"Section {num:02d}")
        self.notes.append(f"Section {num}: {title}.")

    def agenda(self, items):
        s = self.slide(CREAM)
        self.maze(s, "tr")
        self.title(s, "00 · cadrage", "4 heures, un fil très guidé.", "Démonstration lente, exercice court, correction immédiate.")
        y = 2.05
        for i, (label, text, time) in enumerate(items, 1):
            self.text(s, 0.8, y, 0.45, 0.22, f"{i:02d}", 10, CORAL, True)
            self.text(s, 1.35, y, 3.0, 0.22, label, 13, NAVY, True)
            self.text(s, 4.55, y, 5.8, 0.22, text, 11, INK_SOFT)
            self.text(s, 11.2, y, 0.7, 0.22, time, 10, NAVY, True, PP_ALIGN.RIGHT)
            y += 0.48
        self.foot(s, "Sommaire")

    def big_statement(self, crumb, title, subtitle, accent=None):
        s = self.slide(CREAM)
        self.maze(s, "tr")
        self.text(s, 0.55, 0.48, 2.9, 0.16, crumb.upper(), 6, INK_SOFT, True)
        self.text(s, 0.7, 1.8, 10.6, 1.55, title, 42, NAVY, True)
        if accent:
            bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.72), Inches(4.25), Inches(10.4), Inches(0.45))
            bar.fill.solid(); bar.fill.fore_color.rgb = PINK; bar.line.fill.background()
            self.text(s, 0.95, 4.39, 9.7, 0.13, accent, 8, CORAL, True)
        self.text(s, 0.72, 5.05, 8.8, 0.55, subtitle, 15, INK_SOFT)
        self.foot(s, crumb)

    def cards(self, crumb, title, cards, cols=3):
        s = self.slide(CREAM)
        self.title(s, crumb, title)
        start_x = 0.65
        w = (12.0 - (cols - 1) * 0.28) / cols
        for i, (kicker, head, body) in enumerate(cards):
            x = start_x + (i % cols) * (w + 0.28)
            y = 2.25 + (i // cols) * 1.78
            card = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(1.36))
            card.fill.solid(); card.fill.fore_color.rgb = CREAM_SOFT; card.line.color.rgb = RGBColor(210, 207, 198)
            self.text(s, x + 0.18, y + 0.18, w - 0.36, 0.15, kicker.upper(), 6, CORAL, True)
            self.text(s, x + 0.18, y + 0.43, w - 0.36, 0.28, head, 14, NAVY, True)
            self.text(s, x + 0.18, y + 0.82, w - 0.36, 0.38, body, 8.5, INK_SOFT)
        self.foot(s, crumb)

    def activity(self, crumb, title, duration, goal, steps, output):
        s = self.slide(CREAM)
        self.title(s, crumb, title, f"Activité guidée · {duration}")
        badge = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.85), Inches(2.25), Inches(1.15), Inches(1.15))
        badge.fill.solid(); badge.fill.fore_color.rgb = PINK; badge.line.fill.background()
        self.text(s, 1.1, 2.58, 0.65, 0.18, duration, 13, NAVY, True, PP_ALIGN.CENTER)
        self.text(s, 2.45, 2.18, 8.6, 0.55, goal, 24, NAVY, True)
        y = 3.35
        for i, step in enumerate(steps, 1):
            self.text(s, 2.48, y, 0.35, 0.18, f"{i:02d}", 8, CORAL, True)
            self.text(s, 2.95, y, 7.8, 0.22, step, 13, NAVY if i == 1 else INK_SOFT, i == 1)
            y += 0.48
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.45), Inches(5.8), Inches(8.8), Inches(0.55))
        bar.fill.solid(); bar.fill.fore_color.rgb = GREEN; bar.line.fill.background()
        self.text(s, 2.72, 5.99, 8.1, 0.14, f"Sortie attendue · {output}", 8, NAVY, True)
        self.foot(s, crumb)

    def sas_python(self, crumb, title, sas, python, takeaway):
        s = self.slide(CREAM)
        self.title(s, crumb, title)
        left = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(2.1), Inches(5.45), Inches(3.65))
        left.fill.solid(); left.fill.fore_color.rgb = CREAM_SOFT; left.line.fill.background()
        self.text(s, 1.05, 2.4, 4.8, 0.15, "SAS", 8, CORAL, True)
        self.text(s, 1.05, 2.78, 4.8, 2.25, sas, 13, NAVY, False, font="Menlo")
        right = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.55), Inches(2.1), Inches(5.75), Inches(3.65))
        right.fill.solid(); right.fill.fore_color.rgb = BLACK; right.line.fill.background()
        self.text(s, 6.85, 2.4, 4.9, 0.15, "PYTHON / DATABRICKS", 8, CORAL, True)
        self.text(s, 6.85, 2.78, 5.0, 2.25, python, 13, WHITE, False, font="Menlo")
        self.text(s, 0.85, 6.18, 10.8, 0.28, takeaway, 13, NAVY, True)
        self.foot(s, crumb)

    def result_mock(self, crumb, title, rows, interpretation):
        s = self.slide(CREAM)
        self.title(s, crumb, title)
        x0, y0 = 1.05, 2.35
        headers = ["segment", "avg_charges", "lecture métier"]
        widths = [2.6, 2.1, 5.4]
        for i, h in enumerate(headers):
            self.text(s, x0 + sum(widths[:i]), y0, widths[i] - 0.1, 0.15, h.upper(), 7, CORAL, True)
        y = y0 + 0.45
        for r, row in enumerate(rows):
            fill = GREEN if r == 0 else CREAM
            if r == 0:
                rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x0 - 0.12), Inches(y - 0.08), Inches(sum(widths) + 0.2), Inches(0.46))
                rect.fill.solid(); rect.fill.fore_color.rgb = fill; rect.line.fill.background()
            for i, cell in enumerate(row):
                self.text(s, x0 + sum(widths[:i]), y, widths[i] - 0.15, 0.2, cell, 9, NAVY if i == 0 else INK_SOFT, i == 0)
            y += 0.52
        self.text(s, 1.05, 5.85, 9.8, 0.35, interpretation, 16, NAVY, True)
        self.foot(s, crumb)

    def flow(self, crumb, title, steps, caption=None):
        s = self.slide(CREAM)
        self.title(s, crumb, title, caption)
        y = 3.1
        box_w = 2.15
        for i, step in enumerate(steps):
            x = 0.82 + i * 2.55
            fill = CORAL if i == 1 else CREAM_SOFT
            col = WHITE if i == 1 else NAVY
            b = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(box_w), Inches(0.82))
            b.fill.solid(); b.fill.fore_color.rgb = fill; b.line.color.rgb = fill
            self.text(s, x + 0.12, y + 0.28, box_w - 0.24, 0.18, step, 11, col, True, PP_ALIGN.CENTER)
            if i < len(steps) - 1:
                line = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x + box_w + 0.1), Inches(y + 0.41), Inches(x + 2.45), Inches(y + 0.41))
                line.line.color.rgb = NAVY
                line.line.width = Pt(1)
        self.foot(s, crumb)

    def code(self, crumb, title, code, callout):
        s = self.slide(CREAM)
        self.title(s, crumb, title)
        block = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(2.15), Inches(7.15), Inches(3.65))
        block.fill.solid(); block.fill.fore_color.rgb = BLACK; block.line.fill.background()
        self.text(s, 1.05, 2.48, 6.55, 2.8, code, 15, WHITE, False, font="Menlo")
        c = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.35), Inches(2.15), Inches(3.8), Inches(3.65))
        c.fill.solid(); c.fill.fore_color.rgb = PINK; c.line.fill.background()
        self.text(s, 8.65, 2.55, 3.15, 0.16, "LECTURE GUIDÉE", 7, CORAL, True)
        self.text(s, 8.65, 3.05, 3.0, 1.8, callout, 15, NAVY, True)
        self.foot(s, crumb)

    def table(self, crumb, title, headers, rows, highlight_last=False):
        s = self.slide(CREAM)
        self.title(s, crumb, title)
        x0, y0, w = 0.72, 2.05, 11.8
        col_w = w / len(headers)
        for i, h in enumerate(headers):
            self.text(s, x0 + i * col_w, y0, col_w - 0.15, 0.16, h.upper(), 7, CORAL, True)
        y = y0 + 0.42
        for r, row in enumerate(rows):
            fill = GREEN if highlight_last and r == len(rows) - 1 else CREAM
            if fill != CREAM:
                rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x0 - 0.1), Inches(y - 0.08), Inches(w + 0.15), Inches(0.42))
                rect.fill.solid(); rect.fill.fore_color.rgb = fill; rect.line.fill.background()
            line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x0), Inches(y + 0.34), Inches(w), Inches(0.008))
            line.fill.solid(); line.fill.fore_color.rgb = RGBColor(190, 187, 178); line.line.fill.background()
            for i, cell in enumerate(row):
                self.text(s, x0 + i * col_w, y, col_w - 0.18, 0.22, str(cell), 8.5, NAVY if i == 0 else INK_SOFT, i == 0)
            y += 0.52
        self.foot(s, crumb)

    def persona(self, crumb, quote, name="SAS user", role="Analyste métier"):
        s = self.slide(CREAM)
        circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.95), Inches(1.5), Inches(1.1), Inches(1.1))
        circ.fill.solid(); circ.fill.fore_color.rgb = PINK; circ.line.fill.background()
        self.text(s, 1.32, 1.83, 0.35, 0.2, name[0], 22, NAVY, True, PP_ALIGN.CENTER)
        self.text(s, 0.98, 2.78, 1.0, 0.2, name, 8, NAVY, True, PP_ALIGN.CENTER)
        self.text(s, 0.9, 3.02, 1.2, 0.16, role, 6, INK_SOFT, False, PP_ALIGN.CENTER)
        self.text(s, 3.0, 1.55, 7.5, 1.6, quote, 24, NAVY, False)
        self.foot(s, crumb)

    def timeline(self):
        s = self.slide(CREAM)
        self.title(s, "00 · timing", "Déroulé 4h.", "Objectif: alterner explication, démo et pratique.")
        blocks = [
            ("0:00", "Cadrage", CORAL),
            ("0:20", "Databricks", PURPLE),
            ("1:00", "Python", CORAL),
            ("1:45", "Pause", YELLOW),
            ("1:55", "pandas", PURPLE),
            ("2:55", "Spark", CORAL),
            ("3:25", "SQL + challenge", GREEN),
        ]
        x = 0.75
        for time, label, color in blocks:
            b = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(3.1), Inches(1.55), Inches(0.65))
            b.fill.solid(); b.fill.fore_color.rgb = color; b.line.fill.background()
            self.text(s, x + 0.1, 3.24, 1.35, 0.15, time, 8, NAVY, True, PP_ALIGN.CENTER)
            self.text(s, x + 0.1, 3.95, 1.35, 0.25, label, 10, NAVY, True, PP_ALIGN.CENTER)
            x += 1.75
        self.foot(s, "Timing")

    def save(self):
        OUT.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(OUT)
        header = """# Notes formateur - Week 1

Ce deck est conçu pour 4h uniquement si les slides d'activité sont jouées en entier.
Ordre recommandé: 20 min cadrage, 40 min Databricks, 45 min Python, 10 min pause,
60 min pandas, 30 min Spark, 25 min SQL, 10 min challenge/synthèse.

"""
        NOTES.write_text(header + "\n".join(f"- Slide {i+1:02d}: {n}" for i, n in enumerate(self.notes)), encoding="utf-8")


def build_deck() -> None:
    d = Deck()
    d.cover()
    d.agenda([
        ("Cadrage", "objectif, mapping SAS, environnement", "20 min"),
        ("Databricks", "workspace, cluster, notebook, chemins", "40 min"),
        ("Python", "variables, types, listes, fonctions, indentation", "45 min"),
        ("Pause", "respiration + questions", "10 min"),
        ("pandas", "charger, explorer, filtrer, agréger", "60 min"),
        ("Spark", "mêmes gestes sur Spark DataFrame", "30 min"),
        ("SQL", "vue temporaire, COUNT, AVG, GROUP BY", "25 min"),
        ("Challenge", "segment le plus coûteux + restitution", "10 min"),
    ])
    d.timeline()
    d.persona("00 · public", "Je connais SAS. Je veux retrouver mes repères sans devenir développeur Python.")
    d.big_statement("00 · objectif", "Même logique analytique.\nNouvel environnement.", "On garde les réflexes métier: lire, filtrer, agréger, interpréter.", "SAS program → notebook · SAS dataset → DataFrame · PROC SQL → SQL")
    d.table("00 · mapping", "SAS → Databricks.", ["SAS", "Databricks", "À retenir"], [
        ("SAS Program", "Notebook", "Code + texte + résultats"),
        ("DATA Step", "pandas / Spark", "Transformation de table"),
        ("PROC SQL", "SQL cell", "Langage familier"),
        ("SAS Dataset", "DataFrame", "Lignes + colonnes"),
        ("Library", "Catalog / storage", "Emplacement des données"),
    ], True)
    d.activity("00 · tour de table", "Tour de table orienté usage.", "6 min", "Identifier les repères SAS déjà maîtrisés par le groupe.", [
        "Quel PROC utilisez-vous le plus souvent ?",
        "Quelle opération revient chaque semaine : filtrer, agréger, joindre ?",
        "Quelle partie de Databricks vous inquiète le plus ?",
    ], "3 attentes écrites au tableau")
    d.table("00 · dataset", "Dataset fil rouge.", ["Colonne", "Type", "Question métier"], [
        ("age", "numérique", "âge du bénéficiaire"),
        ("sex", "catégorie", "segment homme/femme"),
        ("bmi", "numérique", "indicateur de corpulence"),
        ("children", "numérique", "nombre d'enfants"),
        ("smoker", "catégorie", "facteur de risque"),
        ("region", "catégorie", "zone géographique"),
        ("charges", "numérique", "coût médical annuel"),
    ])
    d.result_mock("00 · question", "Question métier du jour.", [
        ("smoker=yes", "32 100", "segment très coûteux"),
        ("smoker=no", "8 450", "base de comparaison"),
        ("all", "13 700", "moyenne globale"),
    ], "Objectif final: produire ce type de lecture, pas seulement du code.")

    d.divider(1, "Databricks.", "Se repérer avant de coder.", "Section 01 · 35 min")
    d.flow("01 · workspace", "Le parcours d'une analyse.", ["Data", "Notebook", "Cluster", "Result"], "Quatre objets visibles, pas d'architecture complexe.")
    d.cards("01 · repères", "Les zones à reconnaître.", [
        ("Workspace", "Dossiers", "/Shared/training et notebooks"),
        ("Compute", "Cluster", "Le moteur d'exécution"),
        ("Data", "Fichiers", "insurance.csv dans FileStore"),
        ("Notebook", "Cellules", "Markdown, Python, SQL"),
        ("Run", "Résultat", "Sortie sous la cellule"),
        ("Share", "Collaboration", "Même support pour tous"),
    ])
    d.big_statement("01 · démo", "Un notebook s'exécute\ncellule par cellule.", "Ne pas tout lancer au début: on observe le résultat après chaque étape.", "DÉMO LIVE · ouvrir · attacher cluster · exécuter")
    d.activity("01 · démo guidée", "Démo: premier notebook.", "10 min", "Chaque participant voit le même résultat à l'écran.", [
        "Ouvrir /Shared/training/01_Intro_Databricks.",
        "Attacher le cluster single-node.",
        "Exécuter la cellule print puis la cellule variable.",
        "Demander: où apparaît le résultat ?",
    ], "Tout le groupe a exécuté Hello Databricks")
    d.table("01 · chemins", "Deux chemins pour un même fichier.", ["Usage", "Chemin", "Pourquoi"], [
        ("pandas", "/dbfs/FileStore/tables/insurance.csv", "Lecture locale"),
        ("Spark", "/FileStore/tables/insurance.csv", "Lecture Spark"),
        ("Databricks UI", "FileStore/tables/insurance.csv", "Upload visible"),
    ])
    d.activity("01 · exercice", "Exercice: modifier une cellule.", "8 min", "Désacraliser le notebook: modifier, exécuter, observer.", [
        "Changer region = \"North\" en region = \"South\".",
        "Changer premium = 1000 en premium = 1500.",
        "Ré-exécuter seulement la cellule concernée.",
        "Comparer avec son voisin.",
    ], "Une valeur modifiée et un résultat compris")
    d.cards("01 · erreurs", "Erreurs fréquentes Databricks.", [
        ("Cluster", "Non attaché", "La cellule reste en attente ou échoue."),
        ("Ordre", "Cellule oubliée", "La variable n'existe pas encore."),
        ("Chemin", "/dbfs manquant", "pandas ne trouve pas le fichier."),
        ("Run all", "Trop tôt", "On perd le contrôle pédagogique."),
        ("Dossier", "Mauvais import", "Le notebook n'est pas au bon endroit."),
        ("Résultat", "Non lu", "On exécute sans interpréter."),
    ])

    d.divider(2, "Python basics.", "Lire du code analytique simple.", "Section 02 · 50 min")
    d.sas_python("02 · repère", "DATA step → Python.", "data work.demo;\n  age = 45;\n  region = 'southwest';\nrun;", 'age = 45\nregion = "southwest"', "Même idée: nommer des valeurs. Python est plus direct, moins structuré autour d'une table.")
    d.code("02 · variable", "Variable.", 'age = 45\nregion = "southwest"\npremium = 1200.50', "Le signe égal stocke une valeur dans un nom.")
    d.table("02 · types", "4 types à reconnaître.", ["Type", "Exemple", "Usage"], [
        ("int", "45", "âge, compteur"),
        ("float", "27.5", "BMI, montant"),
        ("str", '"southwest"', "modalité texte"),
        ("bool", "True", "condition vraie/fausse"),
    ])
    d.code("02 · liste", "Liste.", 'regions = ["northeast", "northwest",\n           "southeast", "southwest"]\n\nregions[0]', "Une liste regroupe plusieurs valeurs possibles.")
    d.activity("02 · exercice", "Exercice: types et valeurs.", "7 min", "Lire les valeurs avant de parler de syntaxe.", [
        "Modifier age, bmi, smoker.",
        "Exécuter print(type(age)), print(type(bmi)).",
        "Dire à voix haute: entier, décimal, texte ou booléen.",
    ], "Chaque participant sait identifier 3 types")
    d.code("02 · condition", "Condition.", 'if age > 50:\n    print("senior")\nelse:\n    print("non senior")', "L'indentation indique les lignes contrôlées par le if.")
    d.sas_python("02 · condition", "IF SAS → IF Python.", "if age > 50 then segment='senior';\nelse segment='non senior';", 'if age > 50:\n    segment = "senior"\nelse:\n    segment = "non senior"', "La logique est identique. Le vrai changement est l'indentation.")
    d.activity("02 · binôme", "Lecture en binôme.", "6 min", "Forcer une lecture ligne par ligne, sans exécuter d'abord.", [
        "Une personne lit la condition.",
        "L'autre prédit le résultat pour age = 54.",
        "On change age = 35 et on prédit à nouveau.",
    ], "Deux prédictions vérifiées par exécution")
    d.code("02 · fonction", "Fonction.", "def annual_cost(monthly_cost):\n    return monthly_cost * 12\n\nannual_cost(120)", "Une règle métier nommée, réutilisable.")
    d.cards("02 · exercice", "Atelier Python guidé.", [
        ("Étape 1", "Modifier", "Changer age, region, premium"),
        ("Étape 2", "Exécuter", "Observer chaque sortie"),
        ("Étape 3", "Expliquer", "Lire la ligne avec ses mots"),
    ])
    d.activity("02 · correction", "Correction Python.", "8 min", "Ne pas corriger seulement le résultat: corriger la lecture.", [
        "Demander à un participant de lire son code.",
        "Pointer le nom de variable, le signe, la valeur.",
        "Faire expliquer l'indentation par le groupe.",
    ], "Le groupe sait expliquer une cellule courte")
    d.persona("02 · rassurer", "Si je sais lire un IF SAS, je peux lire un IF Python. Le symbole change, la logique reste.")

    d.divider(3, "pandas.", "Manipuler une table comme un analyste.", "Section 03 · 60 min")
    d.big_statement("03 · dataframe", "Un DataFrame est une table.", "Lignes, colonnes, types, aperçu, statistiques.", "SAS dataset · Excel sheet · pandas DataFrame")
    d.table("03 · anatomy", "Anatomie d'une table.", ["Élément", "Dans insurance", "Réflexe"], [
        ("Ligne", "un bénéficiaire", "observation SAS"),
        ("Colonne", "age, bmi, charges", "variable SAS"),
        ("Modalité", "smoker=yes", "valeur catégorielle"),
        ("Mesure", "charges", "variable à analyser"),
    ])
    d.code("03 · load", "Charger le CSV.", 'import pandas as pd\n\ndf = pd.read_csv(\n    "/dbfs/FileStore/tables/insurance.csv"\n)', "Une ligne crée la table df.")
    d.activity("03 · démo", "Démo: chargement contrôlé.", "6 min", "Montrer que charger une table est une étape observable.", [
        "Exécuter import pandas.",
        "Exécuter read_csv.",
        "Afficher type(df).",
        "Afficher len(df).",
    ], "Le groupe voit qu'un DataFrame existe")
    d.flow("03 · workflow", "Workflow pandas.", ["Load", "Inspect", "Filter", "Group", "Interpret"], "Chaque notebook répète ce cycle.")
    d.code("03 · inspect", "Explorer.", "df.head()\ndf.info()\ndf.describe()", "Trois réflexes avant de conclure.")
    d.result_mock("03 · inspect", "Ce qu'on cherche dans l'aperçu.", [
        ("age", "18-64", "ordre de grandeur plausible"),
        ("smoker", "yes/no", "modalité attendue"),
        ("charges", "1 200-55 000", "montants à analyser"),
    ], "Avant de filtrer, on vérifie que la table ressemble à ce qu'on croit.")
    d.activity("03 · exercice", "Exercice: explorer.", "8 min", "Faire produire les premiers contrôles par les participants.", [
        "Afficher les 5 premières lignes.",
        "Compter le nombre de lignes.",
        "Lister les colonnes.",
        "Identifier la colonne cible.",
    ], "Nombre de lignes + liste de colonnes")
    d.code("03 · select", "Sélectionner.", 'df[["age", "charges"]]', "On garde seulement les colonnes utiles.")
    d.sas_python("03 · select", "KEEP SAS → sélection pandas.", "data work.small;\n  set insurance(keep=age charges);\nrun;", 'df[["age", "charges"]]', "Même intention: réduire la table aux variables utiles.")
    d.code("03 · filter", "Filtrer.", 'df[df["smoker"] == "yes"]\n\ndf[df["bmi"] > 30]', "La condition garde certaines lignes.")
    d.sas_python("03 · where", "WHERE SAS → filtre pandas.", "data work.smokers;\n  set insurance;\n  where smoker = 'yes';\nrun;", 'df[df["smoker"] == "yes"]', "Le filtre est une question vraie/fausse posée à chaque ligne.")
    d.activity("03 · exercice", "Exercice: filtrer.", "10 min", "Modifier des seuils sans page blanche.", [
        "Patients age > 50.",
        "Smokers only.",
        "BMI > 30.",
        "Female smokers.",
    ], "4 filtres exécutés et un résultat lu")
    d.code("03 · groupby", "Agréger.", 'df.groupby("region")["charges"].mean()\n\ndf.groupby("smoker")["charges"].mean()', "PROC MEANS mental model: class puis moyenne.")
    d.sas_python("03 · means", "PROC MEANS → groupby.", "proc means data=insurance mean;\n  class smoker;\n  var charges;\nrun;", 'df.groupby("smoker")["charges"].mean()', "Class devient groupby. Var devient la colonne entre crochets.")
    d.result_mock("03 · result", "Lire un résultat groupby.", [
        ("smoker=yes", "32 100", "coûts très supérieurs"),
        ("smoker=no", "8 450", "base non-fumeur"),
        ("gap", "x3.8", "écart métier à commenter"),
    ], "La sortie n'est pas la fin: il faut une phrase métier.")
    d.activity("03 · exercice", "Exercice: agréger.", "10 min", "Faire varier la colonne de regroupement.", [
        "Moyenne de charges par region.",
        "Moyenne de charges par smoker.",
        "Moyenne de charges par sex.",
        "Noter le segment le plus élevé.",
    ], "3 agrégations + une interprétation")
    d.table("03 · mapping", "SAS → pandas.", ["SAS", "pandas", "Geste"], [
        ("PROC PRINT", "head()", "prévisualiser"),
        ("PROC CONTENTS", "info()", "structure"),
        ("WHERE", "df[condition]", "filtrer"),
        ("PROC MEANS", "groupby().mean()", "agréger"),
        ("PROC SORT", "sort_values()", "trier"),
    ], True)
    d.cards("03 · exercice", "Atelier pandas guidé.", [
        ("Explorer", "5 premières lignes", "head, columns, describe"),
        ("Filtrer", "Segments", "age > 50, smoker, BMI"),
        ("Agréger", "Coûts moyens", "region, smoker, sex"),
    ])
    d.cards("03 · pièges", "Pièges pandas du premier jour.", [
        ("Guillemets", "Nom de colonne", 'df["age"], pas df[age]'),
        ("Double crochets", "Plusieurs colonnes", 'df[["age", "charges"]]'),
        ("Égalité", "Deux signes", 'smoker == "yes"'),
        ("Parenthèses", "Filtres combinés", "Chaque condition est isolée"),
        ("Tri", "ascending=False", "Du plus élevé au plus faible"),
        ("Interpréter", "Pas juste afficher", "Lire le résultat métier"),
    ])
    d.big_statement("03 · pause", "Pause 10 minutes.", "Au retour: mêmes gestes, mais avec Spark.", "1:25 → 1:35")

    d.divider(4, "Spark.", "Même table, exécution Databricks.", "Section 04 · 30 min")
    d.flow("04 · mental model", "Spark en une image.", ["Notebook", "Cluster", "Spark DataFrame", "display"], "On évite les détails distribués cette semaine.")
    d.big_statement("04 · message", "Spark n'est pas une nouvelle question.", "C'est une autre manière d'exécuter des gestes table dans Databricks.", "select · filter · groupBy · display")
    d.code("04 · read", "Lire avec Spark.", 'spark_df = spark.read.csv(\n    "/FileStore/tables/insurance.csv",\n    header=True,\n    inferSchema=True\n)', "Même CSV, autre moteur.")
    d.code("04 · transform", "Sélectionner et filtrer.", 'spark_df.select("age", "charges")\n\nspark_df.filter(spark_df.age > 50)', "Même intention que pandas.")
    d.sas_python("04 · translate", "pandas → Spark.", 'df[["age", "charges"]]\n\ndf[df["age"] > 50]', 'spark_df.select("age", "charges")\n\nspark_df.filter(spark_df.age > 50)', "Les verbes deviennent explicites: select, filter.")
    d.activity("04 · exercice", "Exercice: traduire.", "8 min", "Répéter les gestes pandas en Spark.", [
        "select age, charges.",
        "filter age > 50.",
        "display le résultat.",
        "Comparer avec pandas.",
    ], "2 opérations Spark visibles")
    d.code("04 · groupby", "GroupBy Spark.", 'from pyspark.sql.functions import avg\n\nspark_df.groupBy("region").agg(\n    avg("charges").alias("avg_charges")\n)', "La syntaxe change, la question métier reste.")
    d.activity("04 · correction", "Correction Spark.", "7 min", "Insister sur l'intention plus que sur la syntaxe.", [
        "Identifier le verbe Spark.",
        "Identifier la colonne de segment.",
        "Identifier la mesure agrégée.",
    ], "La traduction pandas → Spark est comprise")
    d.table("04 · compare", "pandas ou Spark ?", ["Situation", "pandas", "Spark"], [
        ("Apprendre", "Excellent", "Bien"),
        ("Petit CSV", "Simple", "Possible"),
        ("Databricks", "OK", "Naturel"),
        ("Très gros volume", "Limité", "Prévu pour ça"),
    ], True)

    d.divider(5, "SQL.", "Réutiliser un langage familier.", "Section 05 · 35 min")
    d.code("05 · view", "Créer une vue SQL.", 'spark_df.createOrReplaceTempView("insurance")', "SQL interroge la vue insurance.")
    d.flow("05 · sql flow", "Pourquoi créer une vue ?", ["Spark DataFrame", "Temp view", "SQL cell", "Result"], "La vue fait le pont entre Spark et SQL.")
    d.code("05 · count", "COUNT.", "SELECT COUNT(*) AS row_count\nFROM insurance", "Premier contrôle: le volume chargé.")
    d.code("05 · avg", "GROUP BY.", "SELECT region,\n       AVG(charges) AS avg_charges\nFROM insurance\nGROUP BY region\nORDER BY avg_charges DESC", "PROC SQL mental model.")
    d.sas_python("05 · proc sql", "PROC SQL → Databricks SQL.", "proc sql;\nselect region, mean(charges)\nfrom insurance\ngroup by region;\nquit;", "SELECT region,\n       AVG(charges)\nFROM insurance\nGROUP BY region", "C'est probablement la transition la plus confortable pour un public SAS.")
    d.table("05 · exercices", "Templates SQL.", ["Question", "Mot-clé", "À modifier"], [
        ("Combien de lignes ?", "COUNT", "nom de table"),
        ("Moyenne des charges ?", "AVG", "colonne numérique"),
        ("Par région ?", "GROUP BY", "colonne segment"),
        ("Fumeurs uniquement ?", "WHERE", "condition"),
    ])
    d.activity("05 · exercice", "Exercice SQL.", "8 min", "Compléter des requêtes, pas partir d'une page blanche.", [
        "Changer GROUP BY region en GROUP BY smoker.",
        "Ajouter WHERE age > 50.",
        "Trier avec ORDER BY avg_charges DESC.",
    ], "Une requête SQL modifiée et interprétée")
    d.result_mock("05 · result", "Lire la table SQL.", [
        ("southeast", "14 850", "région la plus élevée"),
        ("northeast", "13 950", "écart modéré"),
        ("northwest", "12 700", "segment plus bas"),
    ], "On commente l'ordre de grandeur, pas seulement le classement.")

    d.divider(6, "Challenge.", "Identifier le segment le plus coûteux.", "Section 06 · 20 min")
    d.flow("06 · méthode", "Méthode attendue.", ["Choose segment", "Filter", "Group", "Sort", "Explain"], "Le livrable est une table et une phrase métier.")
    d.activity("06 · challenge", "Mini challenge.", "12 min", "Répondre à une question métier avec le chemin de son choix.", [
        "Choisir pandas, Spark ou SQL.",
        "Créer un segment: smoker, region, sex ou age_segment.",
        "Calculer la moyenne de charges.",
        "Écrire une phrase d'interprétation.",
    ], "Table triée + phrase métier")
    d.table("06 · grille", "Grille de restitution.", ["Élément", "Attendu", "Exemple"], [
        ("Segment", "variable utilisée", "smoker + age_segment"),
        ("Mesure", "moyenne", "AVG(charges)"),
        ("Tri", "descendant", "coût le plus élevé en haut"),
        ("Phrase", "interprétation", "Les fumeurs >50 ans ont les charges moyennes les plus élevées."),
    ], True)
    d.activity("06 · restitution", "Restitution rapide.", "8 min", "Transformer l'exercice en apprentissage collectif.", [
        "Deux binômes partagent leur segment.",
        "Le groupe vérifie la mesure utilisée.",
        "On reformule une phrase métier propre.",
    ], "Une conclusion validée collectivement")
    d.big_statement("06 · synthèse", "Vous savez ouvrir,\nexécuter,\nmanipuler,\ninterpréter.", "La semaine 2 pourra aller vers Python avancé sans perdre les repères SAS.", "Notebook · Python · pandas · Spark · SQL")
    d.table("06 · checklist", "Checklist de fin de session.", ["Compétence", "Preuve", "OK"], [
        ("Ouvrir un notebook", "01_Intro_Databricks exécuté", "□"),
        ("Modifier du Python", "variable changée", "□"),
        ("Explorer une table", "head/info/describe", "□"),
        ("Filtrer", "smoker ou BMI", "□"),
        ("Agréger", "groupby ou SQL GROUP BY", "□"),
        ("Interpréter", "phrase métier", "□"),
    ])
    d.persona("06 · fin", "Je ne connais pas tout Python, mais je sais exécuter une analyse guidée dans Databricks.", "Learner", "Fin Week 1")
    d.save()


if __name__ == "__main__":
    build_deck()
    print(OUT)
